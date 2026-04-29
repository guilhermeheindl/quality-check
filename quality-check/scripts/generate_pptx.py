"""
Generates Quality Check PPTX from a JSON data file.
Usage: python3 generate_pptx.py <input.json> <output.pptx> <template.pptx>
"""
import sys, json
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from copy import deepcopy
import re

FLAG_EMOJI = {'green': '🟢', 'yellow': '🟡', 'red': '🔴', 'none': '⚪'}

def fmt_brl(val):
    if val is None:
        return '—'
    try:
        return f"R$ {float(val):,.0f}".replace(',', 'X').replace('.', ',').replace('X', '.')
    except Exception:
        return str(val)

def fmt_pct(val):
    if val is None:
        return '—'
    try:
        v = float(val)
        if v <= 1:
            v *= 100
        return f"{v:.1f}%"
    except Exception:
        return str(val)

def set_cell_text(cell, text, bold=False, font_size=None, color=None):
    tf = cell.text_frame
    tf.clear()
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = str(text) if text is not None else ''
    run.font.bold = bold
    if font_size:
        run.font.size = Pt(font_size)
    if color:
        run.font.color.rgb = RGBColor(*color)

def fill_title_slide(slide, data):
    period = data.get('period', '')
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if re.search(r'\d{2}\.\d{2}\.\d{2}', run.text):
                        run.text = period.replace('/', '.')

def fill_overview_slide(slide, data):
    overview = data.get('overview', [])
    gp = data.get('gpName', 'GP')
    obs = data.get('observations', '')
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text
            if 'Overview' in txt or 'GP' in txt:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = run.text.replace('GP Zezinho', f'GP {gp}')
            elif 'Observações' in txt or 'Fatos' in txt:
                tf = shape.text_frame
                tf.clear()
                p0 = tf.paragraphs[0]
                r0 = p0.add_run()
                r0.text = 'Observações Gerais / Fatos Notórios:'
                r0.font.bold = True
                if obs:
                    for line in obs.split('\n'):
                        if line.strip():
                            new_para = tf.add_paragraph()
                            new_run = new_para.add_run()
                            new_run.text = line.strip()
        if hasattr(shape, 'table') and len(overview) > 0:
            tbl = shape.table
            for i, row_data in enumerate(overview):
                row_idx = i + 1
                if row_idx >= len(tbl.rows):
                    break
                row = tbl.rows[row_idx]
                vals = [
                    row_data.get('cliente', ''),
                    row_data.get('fase', ''),
                    FLAG_EMOJI.get(row_data.get('flag', 'none'), '⚪'),
                    str(row_data.get('lt', '') or ''),
                    str(row_data.get('step', '') or ''),
                ]
                for ci, v in enumerate(vals):
                    if ci < len(row.cells):
                        set_cell_text(row.cells[ci], v)

def fill_investimento_slide(slide, data):
    rows = data.get('investimento', [])
    period = data.get('period', '')
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if '/' in run.text and re.search(r'\d', run.text):
                        run.text = f"Investimento e Faturamento / {period}"
        if hasattr(shape, 'table') and len(rows) > 0:
            tbl = shape.table
            for i, row_data in enumerate(rows):
                row_idx = i + 2
                if row_idx >= len(tbl.rows):
                    break
                row = tbl.rows[row_idx]
                pi = row_data.get('percentInvest')
                pf = row_data.get('percentFaturamento')
                flag = 'green' if (pi or 0) >= 0.9 else 'yellow' if (pi or 0) >= 0.7 else 'red'
                vals = [
                    FLAG_EMOJI.get(flag, '⚪'),
                    row_data.get('cliente', ''),
                    fmt_brl(row_data.get('metaInvest')),
                    fmt_brl(row_data.get('atingidoInvestTotal')),
                    fmt_pct(pi),
                    fmt_brl(row_data.get('metaFaturamento')),
                    fmt_brl(row_data.get('atingidoFaturamento')),
                    fmt_pct(pf),
                ]
                for ci, v in enumerate(vals):
                    if ci < len(row.cells):
                        set_cell_text(row.cells[ci], v)

def generate(input_path, output_path, template_path):
    with open(input_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    try:
        prs = Presentation(template_path)
    except Exception:
        prs = Presentation()

    slides = prs.slides
    slide_count = len(slides)

    if slide_count > 0:
        fill_title_slide(slides[0], data)
    if slide_count > 1:
        fill_overview_slide(slides[1], data)
    if slide_count > 2:
        fill_investimento_slide(slides[2], data)

    prs.save(output_path)
    print(f"Saved: {output_path}")

if __name__ == '__main__':
    if len(sys.argv) < 4:
        print("Usage: generate_pptx.py <input.json> <output.pptx> <template.pptx>")
        sys.exit(1)
    generate(sys.argv[1], sys.argv[2], sys.argv[3])
